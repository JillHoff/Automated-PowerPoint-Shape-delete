from pptx import Presentation

prs = Presentation('Outlook Vegas Meats Mar2016.pptx')

for slide in prs.slides:
    for shape in slide.shapes:
        if (shape.name.startswith('Straight Connector')):
            shape.top = 0
            shape.left = 0
            shape.height = 0
            shape.width = 0

prs.save('ppt2.pptx')
